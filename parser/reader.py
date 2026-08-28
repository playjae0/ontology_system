# -*- coding: utf-8 -*-
"""파서 공용 코어 — reader (포맷별 원시 추출)

파서_명세 §3: 포맷 의존은 reader에만 격리한다. doc_type 의존은 adapter에만.
reader는 사람이 1회 작성·유지하는 공용 코드이며 LLM 생성 대상이 아니다.

출력 계약 (adapter의 extract(raw)가 받는 것):

xlsx →
{
  "format": "xlsx",
  "path": "...",
  "sheets": [{
      "name":    str,
      "max_row": int, "max_col": int,
      "cells":   {"A1": value, ...},        # 값이 있는 셀만. 병합은 **전개하지 않음**
      "merged":  ["A4:A31", "B4:B5", ...],  # 병합 범위 원본
      "indent":  {"A5": 3, ...},            # 들여쓰기 수준 (0이 아닌 셀만)
      "bold":    ["A1", "A3", ...],         # 굵은 셀
      "images":  [{"cell": "A17", "ref": "img_001"}, ...],
  }]
}

pptx →
{
  "format": "pptx",
  "path": "...",
  "slides": [{"index": 1, "shapes": ["제목", "본문 …"], "notes": ""}]
}
"""
import logging
import os

_LOG = logging.getLogger("onto.parser.reader")


def read_xlsx(path):
    """**import는 함수 안이다** (문서 7 §7.1 선택 의존의 지연 import 격리).

    최상단 import면 `parser.reader`를 import하는 것만으로 openpyxl이 필요해진다 —
    그러면 패키지 미설치 환경에서 `USE_MOCK=1` 전체 실행이 ImportError로 죽어
    "외부 의존 없이 전체가 동작한다"(문서 1 B12)가 **실측으로** 깨진다. 문면상
    "요구하지 않는다"까지만 두면 최상단 import가 위반이 아니게 되는 것이 그 구멍이다.
    같은 파일의 `read_pptx`가 이미 이 형태다.
    """
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    wb = load_workbook(path, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        cells, indent, bold = {}, {}, []
        for row in ws.iter_rows():
            for c in row:
                if c.value is None:
                    continue
                addr = f"{get_column_letter(c.column)}{c.row}"
                cells[addr] = c.value
                ind = getattr(c.alignment, "indent", 0) or 0
                if ind:
                    indent[addr] = int(ind)
                if c.font and c.font.bold:
                    bold.append(addr)
        images = []
        for i, im in enumerate(getattr(ws, "_images", []), start=1):
            anch = getattr(im, "anchor", None)
            try:
                r = anch._from.row + 1
                col = anch._from.col + 1
                cell = f"{get_column_letter(col)}{r}"
            except Exception:
                cell = None
            images.append({"cell": cell, "ref": f"img_{i:03d}"})
        sheets.append({
            "name": ws.title,
            "max_row": ws.max_row, "max_col": ws.max_column,
            "cells": cells,
            "merged": [str(r) for r in ws.merged_cells.ranges],
            "indent": indent,
            "bold": bold,
            "images": images,
        })
    return {"format": "xlsx", "path": path, "sheets": sheets}


def read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, s in enumerate(prs.slides, start=1):
        shapes = [sh.text_frame.text.strip() for sh in s.shapes
                  if sh.has_text_frame and sh.text_frame.text.strip()]
        notes = ""
        if s.has_notes_slide and s.notes_slide.notes_text_frame is not None:
            notes = s.notes_slide.notes_text_frame.text.strip()
        slides.append({"index": i, "shapes": shapes, "notes": notes})
    return {"format": "pptx", "path": path, "slides": slides}


ENCODINGS = ("utf-8-sig", "cp949", "utf-8")


def _col(i):
    """0-기반 열 번호 → 엑셀 열문자. `openpyxl`을 쓰지 않는다(CSV는 외부 의존 0)."""
    s = ""
    i += 1
    while i:
        i, r = divmod(i - 1, 26)
        s = chr(65 + r) + s
    return s


DELIMS = (",", "\t", ";", "|")


def _delimiter(raw, path):
    """구분자 판정 — **`csv.Sniffer`만 믿지 않는다.** 돌려주는 것은 `(구분자, 근거)`.

    Sniffer는 작은 표에서 실제로 실패한다(실측: 4행 탭 파일에서
    `Could not determine delimiter`). 그때 쉼표로 떨어지면 **탭 파일이 조용히 한 열로
    뭉개지고**, 어댑터는 헤더를 못 찾는데 원인은 화면 어디에도 없다.

    순서: ①`.tsv` 확장자는 탭이다(자명하다) ②Sniffer ③빈도 — **모든 표본 줄에서
    같은 횟수로 1회 이상** 나오는 후보 중 가장 많은 것(표의 열 구분자는 행마다
    같은 수로 나온다) ④쉼표.
    """
    if path.lower().endswith(".tsv"):
        return "\t", "확장자"
    sample = raw[:4096]
    try:
        return _csv_mod().Sniffer().sniff(sample, delimiters="".join(DELIMS)).delimiter, "Sniffer"
    except Exception:
        pass
    lines = [ln for ln in sample.splitlines() if ln.strip()][:10]
    best = None
    for d in DELIMS:
        counts = [ln.count(d) for ln in lines]
        # **모든 줄에 1회 이상** 나와야 후보다 — 값 안에 우연히 섞인 문자를 거른다.
        # 행마다 개수가 같기를 요구하지는 않는다: 열이 덜 찬 짧은 행이 흔하고,
        # 그것까지 요구하면 정상 파일이 후보에서 떨어져 쉼표로 잘못 떨어진다.
        if lines and all(c > 0 for c in counts):
            tot = sum(counts)
            if best is None or tot > best[1]:
                best = (d, tot)
    if best:
        return best[0], "빈도"
    return ",", "기본값"


def _csv_mod():
    import csv
    return csv


def read_csv(path):
    """CSV/TSV → **xlsx와 같은 구조**로 낸다 (어댑터가 포맷을 몰라도 되게).

    **`format`은 `"csv"`로 낸다** — `"xlsx"`로 위장하지 않는다: 어댑터가 포맷별
    분기를 할 수 있어야 하고, 거짓말은 나중에 드러난다. 구조만 같게 해서 xlsx
    어댑터 코드가 거의 그대로 돈다.

    `merged`·`indent`·`bold`·`images`는 **빈 값**이다 — CSV에 그 개념이 없다.
    **키는 둔다**: 어댑터가 참조해도 죽지 않아야 한다.

    **`csv` 표준 모듈에 맡긴다** — 따옴표 안의 쉼표·줄바꿈을 직접 split으로
    다루면 한 셀이 여러 셀로 갈린다.
    """
    import csv as _csv                      # 표준 라이브러리 — 외부 의존 0

    raw, enc, tried = None, None, []
    for e in ENCODINGS:
        tried.append(e)
        try:
            with open(path, encoding=e, newline="") as f:
                raw = f.read()
            enc = e
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if raw is None:
        # **시도한 인코딩을 나열한다** — "못 읽는다"만으로는 다음 수가 안 나온다.
        raise ValueError(f"CSV 인코딩을 판별하지 못했다: {path} — 시도: "
                         f"{', '.join(tried)}. 파일을 UTF-8로 다시 저장하거나 "
                         f"reader.ENCODINGS에 사내 인코딩을 더한다")

    delim, how = _delimiter(raw, path)
    # **추정 결과를 로그로 남긴다** — 탭·세미콜론 파일이 조용히 한 열로 읽히면
    # 어댑터가 헤더를 못 찾고, 그때 원인이 화면 어디에도 없다.
    _LOG.info("CSV %s — 인코딩 %s · 구분자 %r (%s)",
              os.path.basename(path), enc, delim, how)

    rows = list(_csv.reader(raw.splitlines(), delimiter=delim))
    cells, max_col = {}, 0
    for r, row in enumerate(rows, start=1):
        max_col = max(max_col, len(row))     # **최장 행 기준** — 짧은 행은 빈 셀
        for c, v in enumerate(row):
            v = v.strip() if isinstance(v, str) else v
            if v == "" or v is None:
                continue                     # xlsx와 같게 — 값 있는 셀만
            cells[f"{_col(c)}{r}"] = v
    return {"format": "csv", "path": path,
            "sheets": [{"name": os.path.splitext(os.path.basename(path))[0],
                        "max_row": len(rows), "max_col": max_col,
                        "cells": cells,
                        "merged": [], "indent": {}, "bold": [], "images": []}],
            "encoding": enc, "delimiter": delim}


def read(path):
    if path.lower().endswith((".xlsx", ".xlsm")):
        return read_xlsx(path)
    if path.lower().endswith(".pptx"):
        return read_pptx(path)
    if path.lower().endswith((".csv", ".tsv")):
        return read_csv(path)
    raise ValueError(f"지원하지 않는 포맷: {path} — "
                     f"받는 것은 .xlsx · .xlsm · .pptx · .csv · .tsv 다")


# 관찰 범위 — **다단 헤더 문서에서 12줄은 얕다**(B34): 헤더 3행 + 데이터 9행이면
# 병합 블록·값 변형이 2~3회 반복해 나타나는 것을 못 본다(§6.4 관찰 범위 규칙).
# 기본값이 여러 호출부에 흩어져 있었다 — 여기 하나로 모은다.
OBSERVE_ROWS = 20


def head(raw, n=OBSERVE_ROWS):
    """등록 세션에 공급하는 관찰 재료 — head N행 (M5 관찰 범위 규칙)

    **분기는 포맷 이름이 아니라 구조로 한다.** 구판은 `format != "xlsx"`면 슬라이드로
    단정했고, 그래서 `sheets`를 가진 csv가 들어오자 `KeyError: 'slides'`로 죽었다 —
    csv가 xlsx와 같은 구조를 내는 이유가 바로 «어댑터가 포맷을 몰라도 되게»인데,
    이 함수만 이름으로 갈라 그 취지를 깨고 있었다.
    """
    if "sheets" not in raw:
        return {**raw, "slides": raw.get("slides", [])[:n]}
    out = []
    for sh in raw["sheets"]:
        cells = {a: v for a, v in sh["cells"].items()
                 if int("".join(ch for ch in a if ch.isdigit())) <= n}
        # **`max_row`도 함께 줄인다.** 안 줄이면 어댑터가 원래 행 수까지 훑고,
        # 잘린 구간의 빈 셀이 «결측»으로 잡혀 **자르지 않았으면 없었을 C14 실패**가
        # 난다(실측: 병합이 절단면을 걸치면 그 행만 부분적으로 채워진다).
        # 「앞 N행」은 셀만이 아니라 **시트의 크기**까지의 말이다.
        out.append({**sh, "cells": cells,
                    "max_row": min(int(sh.get("max_row") or 0), n),
                    "indent": {a: v for a, v in sh["indent"].items()
                               if int("".join(ch for ch in a if ch.isdigit())) <= n}})
    return {**raw, "sheets": out}
