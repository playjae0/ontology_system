# -*- coding: utf-8 -*-
"""역산 정합 + 시나리오 재료 자동 검증 (v2 — 확대판 대응)

이 스크립트는 두 가지를 한다.
 ① reader + normalizer의 **최소 참조 구현**을 제공한다(병합 전개·상동 해소·복수값 전개·자기완결 검사).
 ② 그 참조 구현으로 raw를 읽어 기존 parsed JSON 명세와 대조한다.

역산 정합의 정의 (가결정 D-18, 확대판에서 정밀화):
 raw를 파싱한 records의 **앞 N건(prefix)**이 mock/parsed/*.json 의 records와 일치한다.
 확대분은 그 뒤에 붙으며 기존 노드·판정에 영향을 주지 않는 신규 항목으로만 구성한다.
"""
import json
import os
import sys
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation

# 레포 루트 기준으로 raw mock 위치를 잡는다 — 이 파일은 tests/ 에 있으므로 부모가 루트다.
# 환경변수 RAW_DIR 로 덮어쓸 수 있다. (구판의 절대경로 하드코딩 대체 — 08-07 13회차 판정)
RAW = os.environ.get("RAW_DIR") or str(Path(__file__).resolve().parent.parent / "mock" / "raw")
DITTO = {"〃", "〝", "same as above", "상동"}

# ---------- 기대값의 정본은 mock/parsed/{doc_id}.json 이다 ----------
# **인라인 기대값을 두지 않는다**(문서 7 §7.5 · impl-B-50). 수기 계약 JSON이 S14의
# 대조 검체인데(§7.5 경로 규약) 스크립트가 그 값을 따로 베껴 들면 **검체가 둘이
# 되고**, 계약이 개정될 때 한쪽만 고쳐진다. 실제로 그 상태였다.
#
# raw의 원본 헤더명과 계약의 필드명은 다르다 — 매핑의 정본은 **어댑터 `expects`의
# `columns`**다(같은 열 문자를 가리키는 두 이름이 거기 나란히 있다). 코드가 매핑을
# 다시 적으면 그것이 세 번째 검체가 된다.
PARSED = Path(__file__).resolve().parent.parent / "mock" / "parsed"
ADAPTERS = Path(__file__).resolve().parent.parent / "mock" / "adapters"


def _adapter(name):
    import importlib.util
    spec = importlib.util.spec_from_file_location(f"_ad_{name}", ADAPTERS / f"{name}.py")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod.ADAPTER


def header_to_field(adapter_name):
    """원본 헤더명 → 계약 필드명. `columns`와 `header_labels`가 **열 문자**로 만난다."""
    exp = _adapter(adapter_name)["expects"]
    cols = exp["columns"]                      # 계약 필드명 → 열 문자
    labels = exp["header_labels"]              # 원본 헤더명 (열 순서)
    letter_of_label = {lab: chr(ord("A") + i) for i, lab in enumerate(labels)}
    out = {}
    for field, letter in cols.items():
        for lab, lt in letter_of_label.items():
            if lt == letter:
                out[lab] = field
    return out


def expected(doc_id, adapter_name, keys, context_key=None):
    """수기 계약 JSON에서 기대 튜플을 만든다 — `keys`는 **원본 헤더명**이다."""
    recs = json.loads((PARSED / f"{doc_id}.json").read_text(encoding="utf-8"))["records"]
    h2f = header_to_field(adapter_name)
    out = []
    for r in recs:
        row = []
        for k in keys:
            field = h2f.get(k, k)
            if field == "context":
                v = (r.get("context") or {}).get(context_key)
            else:
                v = r.get(field)
            row.append(v)
        out.append(tuple(row))
    return out



MULTI_SEP = [",", "/", "\n"]          # 가결정 D-12 기본 닫힌 목록


# ---------- reader + normalizer 참조 구현 ----------
def read_table(path, header_row=3):
    wb = load_workbook(path, data_only=True)
    ws = wb.active
    merged = [str(r) for r in ws.merged_cells.ranges]
    filled = {}
    for rng in ws.merged_cells.ranges:                     # ① 병합 전개
        v = ws.cell(row=rng.min_row, column=rng.min_col).value
        for r in range(rng.min_row, rng.max_row + 1):
            for c in range(rng.min_col, rng.max_col + 1):
                filled[(r, c)] = v
    headers = [ws.cell(row=header_row, column=c).value
               for c in range(1, ws.max_column + 1)]
    headers = [h for h in headers if h is not None]
    rows = []
    for r in range(header_row + 1, ws.max_row + 1):
        vals = [filled.get((r, c), ws.cell(row=r, column=c).value)
                for c in range(1, len(headers) + 1)]
        if all(v is None for v in vals):
            continue
        rows.append((r, dict(zip(headers, vals))))
    return headers, rows, merged


def normalize(headers, rows, entity_cols, required):
    """② 상동 해소 ③ 복수값 전개 + 자기완결 검사"""
    out, failures, prev, dittos = [], [], {}, 0
    for rno, rec in rows:
        rec = dict(rec)
        for h in headers:
            v = rec.get(h)
            if isinstance(v, str) and v.strip() in DITTO:
                rec[h] = prev.get(h)
                dittos += 1
            elif v is not None:
                prev[h] = v
        miss = [c for c in required if not rec.get(c)]
        if miss:
            failures.append((rno, f"필수 필드 결측: {miss}"))
            continue
        expand = []
        for c in entity_cols:
            v = rec.get(c)
            if isinstance(v, str):
                for sep in MULTI_SEP:
                    if sep in v:
                        expand.append((c, [x.strip() for x in v.split(sep) if x.strip()]))
                        break
        if expand:
            col, parts = expand[0]
            for p in parts:
                nr = dict(rec)
                nr[col] = p
                out.append((rno, nr))
        else:
            out.append((rno, rec))
    return out, failures, dittos


def show(label, ok, detail=""):
    print(f"  [{'PASS' if ok else 'FAIL'}] {label}" + (f"  — {detail}" if detail else ""))
    return bool(ok)


def diff_prefix(got, exp, tag):
    for i, (g, e) in enumerate(zip(got, exp), 1):
        if g != e:
            print(f"      {tag}{i}  got={g}\n           exp={e}")


allok = True
CP_COLS = ["공정구분", "공정번호", "공정명", "극성", "설비", "관리항목",
           "규격", "측정방법", "대응계획", "적용모델"]

# ============================================================
print("\n■ CP01.xlsx — 역산 정합 prefix 12건 + 확대분")
h, rows, merged = read_table(f"{RAW}/CP01.xlsx")
recs, fails, dit = normalize(h, rows, ["설비", "관리항목"], ["공정명", "설비", "관리항목"])
# 튜플 = (공정구분, 공정명, 극성, 설비, 관리항목, 규격, 적용모델)
# 공정구분은 C11(coord_mismatch)을 어서션하기 위해 08-07에 추가됐다 — 그 전에는
# 전 행이 "조립"이라 비교에 넣을 이유가 없었다.
CP_KEYS = ["공정구분", "공정명", "극성", "설비", "관리항목", "규격", "적용모델"]
EXP_CP = expected("CP01", "cp", CP_KEYS, context_key="model")
got = [(r["공정구분"], r["공정명"], r["극성"], r["설비"], r["관리항목"],
        r["규격"], r["적용모델"]) for _, r in recs]
allok &= show("헤더 = cp 계약 10열", h == CP_COLS, str(h))
allok &= show("병합 20건 이상", len(merged) >= 20, f"{len(merged)}건")
allok &= show("상동(〃) 5건 이상 해소", dit >= 5, f"{dit}건")
allok &= show("자기완결 실패 0건", not fails, str(fails))
allok &= show(f"record {len(got)}건 (확대 목표 24건 이상)", len(got) >= 24)
allok &= show("prefix 12건이 기존 parsed 명세와 일치", got[:12] == EXP_CP)
if got[:12] != EXP_CP:
    diff_prefix(got[:12], EXP_CP, "C")
# 확대분이 기존 판정을 흔들지 않는지
alig = [g for g in got if g[4] == "적층 정렬도"]
allok &= show("확대분이 '적층 정렬도' spec 판정을 건드리지 않음 (M1 2건·M2 1건 유지)",
              len(alig) == 3, f"{len(alig)}건")
clr = [(g[5], g[6]) for g in got if g[4] == "금형 클리어런스"]
allok &= show("'금형 클리어런스'에 M3 맥락 항목이 병렬 추가됨",
              sorted(clr) == [("20±2㎛", None), ("22±2㎛", "M3")], str(clr))
# 골격 v3.2가 해소해야 하는 raw 표면형 8종. 전부 골격 노드로 착지하되 경로가 다르다 —
# 개념 노드 직행(노칭·스태킹·패키징) / 인스턴스 auto alias(cathode·anode 탭용접) /
# seed ALIASES(전해액주입 → 패키징::전해액 주액) / 짧은 이름 auto alias(사이드·프리 실링).
SKEL = ["노칭", "스태킹", "cathode 탭용접", "anode 탭용접", "패키징", "전해액주입",
        "사이드 실링", "프리 실링"]
used = {g[1] for g in got}
allok &= show("골격 표면형 8종 중 CP01이 커버하는 분량", len(set(SKEL) & used) >= 6,
              f"미사용={sorted(set(SKEL) - used)}")
# "탭용접"은 이제 개념 노드에 해소된다(A11-6 저해상도 부착) — 구 '극성 모호' 대조군이
# 아니라 **개념 해상도 부착**의 대조군이다.
allok &= show("개념 해상도 표면형('탭용접')이 대조군으로 남아 있음", "탭용접" in used)

# ============================================================
print("\n■ PFMEA01.xlsx — 역산 정합 prefix 13건 + 확대분")
h, rows, merged = read_table(f"{RAW}/PFMEA01.xlsx")
recs, fails, dit = normalize(h, rows, ["고장모드", "고장원인"],
                             ["공정명", "고장모드", "고장원인"])
FM_KEYS = ["공정명", "고장모드", "고장원인", "영향분류", "심각도", "관리항목(원인)"]
EXP_FM = expected("PFMEA01", "pfmea", FM_KEYS)
got = [(r["공정명"], r["고장모드"], r["고장원인"], r["영향분류"], r["심각도"],
        r["관리항목(원인)"]) for _, r in recs]
allok &= show("자기완결 실패 0건", not fails, str(fails))
allok &= show(f"record {len(got)}건 (확대 목표 30건 이상)", len(got) >= 30)
allok &= show("prefix 13건이 기존 parsed 명세와 일치", got[:13] == EXP_FM)
if got[:13] != EXP_FM:
    diff_prefix(got[:13], EXP_FM, "R")
allok &= show("nested 병합 다수 (30건 이상)", len(merged) >= 30, f"{len(merged)}건")
allok &= show("R12만 비고 열 값 보유 → unknown_field 재료",
              [r["비고"] for _, r in recs].count("재발 2건") == 1)
SEV = {"단락": 9, "화재": 9, "방전기능상실": 8, "충전기능상실": 7}
allok &= show("severity ↔ effect 1:1 정렬 (골격 내 effect 전량)",
              all(g[4] == SEV[g[3]] for g in got if g[3] in SEV))
_mold = {g[0] for g in got if g[2] == "금형 마모"}
_foreign = {g[0] for g in got if g[5] == "이물 검출 감도"}
allok &= show("'금형 마모'가 2개 공정에 걸쳐 등장 (교차 매칭 재료)", len(_mold) >= 2, str(sorted(_mold)))
allok &= show("'이물 검출 감도'가 3개 공정에 걸쳐 등장 (걸침 Property 교차 매칭)",
              len(_foreign) >= 3, str(sorted(_foreign)))
allok &= show("causes 2단 연쇄 존재 (분리막 주름 → 내부 단락)",
              any(g[1] == "분리막 주름" for g in got)
              and any(g[2] == "분리막 주름" for g in got))
allok &= show("골격 밖 effect('셀 부풀음') 2건 — orphan_anchor 재료",
              sum(1 for g in got if g[3] == "셀 부풀음") == 2)

# ============================================================
print("\n■ CP02_drift.xlsx — preflight 불일치 (S7)")
h, rows, _ = read_table(f"{RAW}/CP02_drift.xlsx")
diff = [(a, b) for a, b in zip(CP_COLS, h) if a != b]
allok &= show("헤더가 cp expects와 불일치(=중단 사유 발생)", bool(diff), str(diff))
allok &= show("불일치가 1개 열로 국소적", len(diff) == 1)
allok &= show("데이터는 정상 (extract 미실행을 확인할 수 있음)", len(rows) >= 14, f"{len(rows)}행")

# ============================================================
print("\n■ CP03_bad.xlsx — 자기완결 불성립 → parse_failure (S8)")
h, rows, _ = read_table(f"{RAW}/CP03_bad.xlsx")
recs, fails, _ = normalize(h, rows, ["설비", "관리항목"], ["공정명", "설비", "관리항목"])
allok &= show("계약 위반 행 정확히 1건 검출", len(fails) == 1, str(fails))
allok &= show("나머지 행은 정상 (문서 단위 실패의 대조군)", len(recs) >= 10, f"{len(recs)}건")

# ============================================================
print("\n■ CP04_unlabeled.xlsx — 지문 스캔(S11) + 복수값 전개 2건")
h, rows, _ = read_table(f"{RAW}/CP04_unlabeled.xlsx")
recs, fails, _ = normalize(h, rows, ["설비", "관리항목"], ["공정명", "설비", "관리항목"])
allok &= show("헤더가 cp expects와 완전 일치 → 후보 'cp' 제안 성립", h == CP_COLS)
allok &= show(f"raw {len(rows)}행 → record {len(recs)}건 (복수값 2건 전개)",
              len(recs) == len(rows) + 2)
items = {r["관리항목"] for _, r in recs}
allok &= show("전개 결과에 4개 항목 각각 존재",
              {"주액량", "주액 속도", "실링 폭", "실링 강도"} <= items)

# ============================================================
print("\n■ IPQC01/IPQC02.xlsx — 신규 정형 doc_type 등록 (role 배정표·UNMAPPABLE)")
IPQC_COLS = ["대공정", "공정No", "공정명", "극성", "검사설비", "검사항목", "규격",
             "측정방법", "판정기준", "부적합 조치", "적용모델", "검사자", "검사일시",
             "성적서번호", "최근 불량 이력", "관련 표준문서"]
h1, r1, m1 = read_table(f"{RAW}/IPQC01.xlsx")
rec1, f1, d1 = normalize(h1, r1, ["검사설비", "검사항목"], ["공정명", "검사설비", "검사항목"])
h2, r2, m2 = read_table(f"{RAW}/IPQC02.xlsx")
rec2, f2, d2 = normalize(h2, r2, ["검사설비", "검사항목"], ["공정명", "검사설비", "검사항목"])
allok &= show("헤더 16열 (유효 11 + meta 3 + UNMAPPABLE 2)",
              h1 == IPQC_COLS and h2 == IPQC_COLS)
allok &= show("자기완결 실패 0건", not f1 and not f2, str(f1 + f2))
allok &= show(f"IPQC01 raw {len(r1)}행 → record {len(rec1)}건 (확대 목표 30건 이상)",
              len(r1) >= 30 and len(rec1) == len(r1) + 1)
allok &= show(f"IPQC02 raw {len(r2)}행 (표본 2부 변형 — 행 수 다름)",
              len(r2) >= 18 and len(r2) != len(r1))
used = {r["공정명"] for _, r in rec1} | {r["공정명"] for _, r in rec2}
allok &= show("골격 표면형 8종 전량 소진", set(SKEL) <= used,
              f"미사용={sorted(set(SKEL) - used)}")
# process_no는 meta라 골격 판정과 무관하다(D-48) — 재매핑에서 재부여하지 않았으므로
# 사이드·프리 실링이 같은 OP-60을 공유한다. 여기서 보는 것은 표본 내부의 일관성뿐이다.
ORDER = {"노칭": "OP-10", "스태킹": "OP-20", "cathode 탭용접": "OP-30",
         "anode 탭용접": "OP-31", "패키징": "OP-40", "전해액주입": "OP-50",
         "사이드 실링": "OP-60", "프리 실링": "OP-60"}
bad = [(r["공정명"], r["공정No"]) for _, r in rec1 + rec2
       if ORDER.get(r["공정명"]) and r["공정No"] != ORDER[r["공정명"]]]
allok &= show("공정번호가 골격 순서에 정렬 (D-24)", not bad, str(bad))
allok &= show("meta 후보 3열 전량 채움",
              all(r["검사자"] and r["검사일시"] and r["성적서번호"] for _, r in rec1))
allok &= show("UNMAPPABLE 후보 '관련 표준문서' 양쪽 전량 채움",
              all(r["관련 표준문서"] for _, r in rec1 + rec2))
allok &= show("UNMAPPABLE 후보 '최근 불량 이력' — 01 부분 채움 / 02 전량 공란 (표본 2부 변형)",
              any(r["최근 불량 이력"] for _, r in rec1)
              and not any(r["최근 불량 이력"] for _, r in rec2))
allok &= show("대공정 채움 전략이 표본마다 다름 (01 병합 / 02 행별 반복)",
              any(str(x).startswith("A4:A") for x in m1)
              and not any(str(x).startswith("A4:A") for x in m2))
allok &= show("맥락형(M2) 항목 다수 — 병렬 저장 재료",
              sum(1 for _, r in rec1 if r["적용모델"] == "M2") >= 3)

# ============================================================
print("\n■ TOC01/TOC02.xlsx — 목차형 3단 계층 (S1·S10)")
import re
def toc_stats(path):
    ws = load_workbook(path).active
    lines = {r: (ws.cell(row=r, column=1).value,
                 ws.cell(row=r, column=1).alignment.indent or 0)
             for r in range(1, ws.max_row + 1) if ws.cell(row=r, column=1).value}
    lv1 = [r for r, (v, i) in lines.items() if r > 1 and i == 0 and re.match(r"^\d+\.\s", str(v))]
    return ws, lines, lv1

ws1, lines1, lv1_1 = toc_stats(f"{RAW}/TOC01.xlsx")
ws2, lines2, lv1_2 = toc_stats(f"{RAW}/TOC02.xlsx")
allok &= show(f"TOC01 본문 {len(lines1)}행 (확대 목표 25행 이상)", len(lines1) >= 25)
allok &= show("레벨1 3회 반복 — 관찰 확장 규칙(M5) 재료 강화", len(lv1_1) >= 3, f"행={lv1_1}")
allok &= show("3단 계층 존재 (indent 0·1·2·3 전부)",
              {0, 1, 2, 3} <= {i for _, i in lines1.values()})
allok &= show("들여쓰기 이상치(5) 존재 → '판정 불가' 유도",
              5 in {i for _, i in lines1.values()})
allok &= show("이미지 2건 → placeholder 조각", len(ws1._images) == 2, f"{len(ws1._images)}건")
allok &= show(f"TOC02 본문 {len(lines2)}행 · 구조가 다름 (레벨3 없음)",
              len(lines2) >= 20 and 2 not in {i for _, i in lines2.values()})
allok &= show("TOC02에는 이상치 없음 (정상 대조군)", 5 not in {i for _, i in lines2.values()})
allok &= show("TOC01·TOC02 행 수가 다름", len(lines1) != len(lines2))

# ============================================================
print("\n■ PPT_basic.pptx — 기본 어댑터 (S9)")
prs = Presentation(f"{RAW}/PPT_basic.pptx")
n = len(prs.slides._sldIdLst)
texts = [" ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame).strip()
         for s in prs.slides]
allok &= show(f"슬라이드 {n}장 = 청크 {n}개 (확대 목표 8장 이상)", n >= 8)
allok &= show("전 슬라이드에 텍스트 존재 (빈 청크 없음)", all(texts))
allok &= show("슬라이드당 평균 40자 이상 (추출 재료로 충분)",
              sum(len(t) for t in texts) / n >= 40, f"평균 {sum(len(t) for t in texts)//n}자")

print("\n" + "=" * 62)
print("전체 결과:", "PASS — 역산 정합 성립 · 확대 목표 충족" if allok else "FAIL — 위 항목 확인 필요")
sys.exit(0 if allok else 1)
